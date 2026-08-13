"""
Source fetch, text extraction, and AEGIS-format index build (PROMPTS 1.1).

The eval host's network is unrestricted, so HTML and PDF documents are
fetched with `urllib` and saved as `.txt` plus an AEGIS-format Jaccard
chunk index (`aegis.llm.rag_retriever.build_rag_index`). The same index
serves the source-conditioned generator (PROMPTS Section 1) and is shaped
to be loadable by the MANDATE-primary Procedure role's retriever at
runtime.

This module deliberately uses the standard library for HTTP, with
`pypdf` as the only optional dependency for PDF text extraction; both are
present in the eval-host venv. The sandbox environment cannot fetch many
of the curated URLs, so the unit tests use file:// URLs and a small
sample PDF; the live eval-host run is driven by Handoff 07.

Every URL that fails is recorded as a failed entry in `BuildReport`; no
text is fabricated.
"""
from __future__ import annotations

import html
import io
import os
import re
import time
import urllib.request
from dataclasses import dataclass, field, asdict
from typing import Optional
from urllib.error import URLError, HTTPError

from .curated_sources import SourceSpec

# Identifying academic User-Agent. The 2026-06-03 Handoff 07d diagnostic
# tested swapping to a generic Chrome 120 string; the change made things
# strictly worse (3 working sources went to 0; even the previously-working
# ODNI ATAs and intelligence.gov began 403'ing). The plausible explanation
# is that .gov bot detection treats a bare Chrome User-Agent without the
# accompanying `sec-ch-ua` and `sec-fetch-*` headers as more suspicious
# than a clearly-identified academic fetcher. The User-Agent is back to
# the prior identifying string; the Accept and Accept-Language headers
# stay in (they were not the cause). See `handoffs/HANDOFF_07d_report_
# 2026-06-03.md` for the evidence.
DEFAULT_USER_AGENT = ("Mozilla/5.0 (MANDATE-eval-2026Q2 corpus fetcher; "
                      "academic use; replication package)")
DEFAULT_ACCEPT = ("text/html,application/xhtml+xml,application/xml;q=0.9,"
                  "application/pdf;q=0.9,*/*;q=0.8")
DEFAULT_ACCEPT_LANGUAGE = "en-US,en;q=0.9"
DEFAULT_TIMEOUT_S = 30.0
DEFAULT_CHUNK_SIZE = 1200
DEFAULT_CHUNK_OVERLAP = 200


@dataclass
class SourceFetchResult:
    spec: SourceSpec
    saved_path: str = ""
    bytes_text: int = 0
    error: str = ""

    @property
    def ok(self) -> bool:
        return not self.error and self.bytes_text > 0

    def to_dict(self) -> dict:
        d = asdict(self.spec)
        d.update({"saved_path": self.saved_path,
                  "bytes_text": self.bytes_text,
                  "error": self.error, "ok": self.ok})
        return d


@dataclass
class BuildReport:
    domain: str
    sources_dir: str
    index_path: str
    fetched: list = field(default_factory=list)
    failed: list = field(default_factory=list)
    files_indexed: int = 0
    chunks_indexed: int = 0

    def to_dict(self) -> dict:
        return {"domain": self.domain, "sources_dir": self.sources_dir,
                "index_path": self.index_path,
                "fetched": [r.to_dict() for r in self.fetched],
                "failed": [r.to_dict() for r in self.failed],
                "files_indexed": self.files_indexed,
                "chunks_indexed": self.chunks_indexed}


# --- text extraction --------------------------------------------------------

def _strip_html(html_text: str) -> str:
    """Compact HTML to readable text. Removes script/style blocks, replaces
    tags with whitespace, decodes entities, collapses whitespace."""
    s = re.sub(r"(?is)<script.*?>.*?</script>", " ", html_text)
    s = re.sub(r"(?is)<style.*?>.*?</style>", " ", s)
    s = re.sub(r"<[^>]+>", " ", s)
    s = html.unescape(s)
    s = re.sub(r"\s+", " ", s)
    return s.strip()


def extract_html_text(raw_bytes: bytes) -> str:
    """Decode HTML bytes and return a plain-text body."""
    for enc in ("utf-8", "latin-1"):
        try:
            return _strip_html(raw_bytes.decode(enc))
        except UnicodeDecodeError:
            continue
    return _strip_html(raw_bytes.decode("utf-8", errors="ignore"))


def extract_pdf_text(raw_bytes: bytes) -> str:
    """Extract text from a PDF byte string using pypdf. Returns the full
    document text concatenated across pages, with a page-break newline."""
    try:
        import pypdf
    except ImportError as e:
        raise ImportError(
            "pypdf is required for PDF text extraction; "
            "pip install pypdf in the eval-host venv") from e
    reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
    pages = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            pages.append("")
    return "\n".join(pages).strip()


def extract_docx_text(raw_bytes: bytes) -> str:
    """Extract text from a .docx byte string using python-docx. Includes
    paragraphs, table cells, and headers/footers."""
    try:
        from docx import Document
    except ImportError as e:
        raise ImportError(
            "python-docx is required for .docx text extraction; "
            "pip install python-docx in the eval-host venv") from e
    doc = Document(io.BytesIO(raw_bytes))
    parts = []
    for p in doc.paragraphs:
        if p.text:
            parts.append(p.text)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                if cell.text:
                    parts.append(cell.text)
    return "\n".join(parts).strip()


def extract_pptx_text(raw_bytes: bytes) -> str:
    """Extract text from a .pptx byte string using python-pptx. Includes
    slide titles, body text, notes pages, and table cells."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError(
            "python-pptx is required for .pptx text extraction; "
            "pip install python-pptx in the eval-host venv") from e
    prs = Presentation(io.BytesIO(raw_bytes))
    parts = []
    for slide_i, slide in enumerate(prs.slides, start=1):
        parts.append("=== Slide %d ===" % slide_i)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text:
                        parts.append(p.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            parts.append(cell.text)
        # speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            parts.append("[notes] " + slide.notes_slide.notes_text_frame.text)
    return "\n".join(parts).strip()


# --- HTTP -------------------------------------------------------------------

def _http_get(url: str, *, user_agent: str = DEFAULT_USER_AGENT,
              timeout: float = DEFAULT_TIMEOUT_S) -> bytes:
    req = urllib.request.Request(url, headers={
        "User-Agent": user_agent,
        "Accept": DEFAULT_ACCEPT,
        "Accept-Language": DEFAULT_ACCEPT_LANGUAGE,
    })
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return resp.read()


def _slugify(name: str) -> str:
    s = re.sub(r"[^A-Za-z0-9_.-]+", "_", name).strip("_") or "source"
    return s[:80]


# --- per-source fetch -------------------------------------------------------

def fetch_one(spec: SourceSpec, *, out_dir: str,
              user_agent: str = DEFAULT_USER_AGENT,
              timeout: float = DEFAULT_TIMEOUT_S,
              http_getter=None) -> SourceFetchResult:
    """Fetch one source, extract text, save to `out_dir/<slug>.txt`.

    `http_getter` is an injection point for tests: a callable that takes
    `(url, timeout=...)` and returns the raw bytes; defaults to urllib.
    """
    res = SourceFetchResult(spec=spec)
    os.makedirs(out_dir, exist_ok=True)
    slug = _slugify(spec.title)
    out_path = os.path.join(out_dir, "%s.txt" % slug)
    getter = http_getter or (lambda u, timeout=timeout:
                              _http_get(u, user_agent=user_agent,
                                          timeout=timeout))
    try:
        raw = getter(spec.url)
    except (URLError, HTTPError) as e:
        res.error = "fetch failed: %r" % e
        return res
    except Exception as e:
        res.error = "fetch failed: %r" % e
        return res

    try:
        if spec.media == "pdf":
            text = extract_pdf_text(raw)
        elif spec.media == "html":
            text = extract_html_text(raw)
        elif spec.media == "docx":
            text = extract_docx_text(raw)
        elif spec.media == "pptx":
            text = extract_pptx_text(raw)
        else:
            res.error = "unknown media: %r" % spec.media
            return res
    except Exception as e:
        res.error = "extract failed: %r" % e
        return res

    text = text.strip()
    if not text:
        res.error = "extracted text was empty"
        return res

    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    res.saved_path = out_path
    res.bytes_text = len(text.encode("utf-8"))
    return res


def fetch_all(specs: list, *, out_dir: str,
              http_getter=None,
              sleep_between_s: float = 0.5) -> tuple:
    """Fetch every spec; return (fetched_results, failed_results). A small
    delay between requests is polite to public hosts."""
    fetched, failed = [], []
    for spec in specs:
        res = fetch_one(spec, out_dir=out_dir, http_getter=http_getter)
        if res.ok:
            fetched.append(res)
        else:
            failed.append(res)
        if sleep_between_s > 0 and not http_getter:
            time.sleep(sleep_between_s)
    return fetched, failed


# --- index build (AEGIS Jaccard format) -------------------------------------

def build_domain_index(*, domain: str, project_root: str,
                       aegis_eval_src: Optional[str] = None,
                       http_getter=None,
                       sleep_between_s: float = 0.5) -> BuildReport:
    """End-to-end build for one domain: fetch the curated sources, save
    text, and produce the AEGIS-format Jaccard index. Returns a
    BuildReport. The index lands at `rag/embeddings/<domain>.jsonl` under
    the project root.
    """
    from .curated_sources import get_sources
    aegis_src = aegis_eval_src or os.path.join(project_root,
                                                 "AEGIS-eval", "src")
    if aegis_src not in os.sys.path:
        os.sys.path.insert(0, aegis_src)
    from aegis.llm.rag_retriever import build_rag_index

    out_dir = os.path.join(project_root, "rag", "sources", domain)
    index_path = os.path.join(project_root, "rag", "embeddings",
                              "%s.jsonl" % domain)
    os.makedirs(out_dir, exist_ok=True)
    os.makedirs(os.path.dirname(index_path), exist_ok=True)

    specs = get_sources(domain)
    fetched, failed = fetch_all(specs, out_dir=out_dir,
                                http_getter=http_getter,
                                sleep_between_s=sleep_between_s)

    summary = build_rag_index(input_dir=out_dir, output_path=index_path,
                               source=domain.upper(),
                               chunk_size=DEFAULT_CHUNK_SIZE,
                               chunk_overlap=DEFAULT_CHUNK_OVERLAP)
    return BuildReport(
        domain=domain, sources_dir=out_dir, index_path=index_path,
        fetched=fetched, failed=failed,
        files_indexed=summary.get("files_indexed", 0),
        chunks_indexed=summary.get("chunks_indexed", 0))
